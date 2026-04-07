#!/usr/bin/env python3
"""
Generate PPT courseware for Dora-MoveIt2 + GEN72 dual-arm course.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MEDIUM = RGBColor(0x16, 0x21, 0x3E)
BG_CARD = RGBColor(0x1F, 0x2B, 0x4D)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xDD, 0xEE)
TEXT_DIM = RGBColor(0x88, 0x99, 0xAA)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)
CODE_GREEN = RGBColor(0x7E, 0xE7, 0x87)
CODE_YELLOW = RGBColor(0xE6, 0xDB, 0x74)
CODE_CYAN = RGBColor(0x66, 0xD9, 0xEF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_LIGHT, bullet_color=ACCENT_BLUE, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        # bullet marker
        run_b = p.add_run()
        run_b.text = "\u25B8 "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = font_name

        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = font_name
    return txBox


def add_code_block(slide, left, top, width, height, code, font_size=12):
    shape = add_shape_bg(slide, left, top, width, height, CODE_BG, radius=0.02)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_right = Pt(12)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(code.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = CODE_GREEN
        run.font.name = "Consolas"
    return shape


def add_table_slide(slide, left, top, col_widths, headers, rows,
                    header_color=ACCENT_BLUE, font_size=13):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = sum(col_widths)
    height = Inches(0.35) * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # header
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    # rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size - 1)
                p.font.color.rgb = TEXT_LIGHT
                p.font.name = "Microsoft YaHei"
            bg = BG_CARD if ri % 2 == 0 else BG_MEDIUM
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    return table_shape


def add_chapter_header(slide, chapter_num, title_cn, title_en):
    set_slide_bg(slide, BG_DARK)
    # chapter number badge
    shape = add_shape_bg(slide, Inches(0.8), Inches(0.6), Inches(1.4), Inches(0.5), ACCENT_BLUE, 0.5)
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = f"Chapter {chapter_num}"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_WHITE
    run.font.name = "Consolas"

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.8),
                 title_cn, font_size=36, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(11), Inches(0.5),
                 title_en, font_size=20, color=TEXT_DIM, bold=False)


def add_section_title(slide, title, subtitle=None):
    set_slide_bg(slide, BG_DARK)
    add_shape_bg(slide, Inches(0.8), Inches(2.5), Inches(0.08), Inches(0.6), ACCENT_BLUE)
    add_text_box(slide, Inches(1.1), Inches(2.3), Inches(11), Inches(1.0),
                 title, font_size=28, color=TEXT_WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1.1), Inches(3.3), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=TEXT_DIM)


# ══════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ── TITLE SLIDE ──
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Dora-MoveIt2 + GEN72\n双臂机器人精简课程",
             font_size=44, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "基于 dora-rs 数据流框架 · MuJoCo 仿真 · Python 全栈",
             font_size=20, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "纯Python · pip安装 · 5分钟上手 · API兼容ROS MoveIt",
             font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
             "目标硬件：双 Realman GEN72 7-DOF 机械臂",
             font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)


# ── COURSE OVERVIEW ──
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "课程大纲", font_size=32, color=TEXT_WHITE, bold=True)
add_table_slide(slide, Inches(0.8), Inches(1.2),
    [Inches(1.2), Inches(4.5), Inches(5.5)],
    ["章节", "主题", "关键内容"],
    [
        ["第1章", "环境搭建", "pip install, MuJoCo, dora up"],
        ["第2章", "数据流基础", "操作器, YAML数据流, Config注入"],
        ["第3章", "GEN72模型", "Python Config类, 关节限制, FK链"],
        ["第4章", "配置体系", "RobotConfig协议, DualArmConfig"],
        ["第5章", "MoveGroup API", "go(), set_named_target(), 笛卡尔路径"],
        ["第6章", "单臂规划", "关节目标, 笛卡尔路径, 抓取, 避障"],
        ["第7章", "双臂协调", "DualMoveGroup, 14D规划, 同步运动"],
        ["第8章", "综合实验", "抓取-传递-放置 完整Demo"],
    ])


# ══════════════════════════════════════════════════
#  CHAPTER 1: Environment Setup
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_chapter_header(slide, 1, "课程介绍与环境搭建", "Course Introduction & Environment Setup")

# 1 - Why Dora
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "为什么选择 Dora-MoveIt2？")

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "ROS2 + MoveIt2 的痛点", font_size=28, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(2.5), [
    "安装复杂：需特定Ubuntu版本，C++编译30分钟+",
    "概念繁多：节点/话题/服务/TF树/URDF/SRDF...",
    "调试困难：DDS中间件 + colcon构建系统",
], bullet_color=ACCENT_RED)

add_text_box(slide, Inches(6.8), Inches(0.4), Inches(6), Inches(0.6),
             "Dora-MoveIt2 的优势", font_size=28, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(2.5), [
    "纯Python：pip install 即装即用",
    "数据流YAML：一个文件定义所有连接",
    "MuJoCo仿真：自带3D可视化",
    "API一致：与ROS MoveIt几乎完全一致",
], bullet_color=ACCENT_GREEN)

add_table_slide(slide, Inches(0.8), Inches(4.0),
    [Inches(3.0), Inches(4.0), Inches(4.0)],
    ["对比项", "ROS2 + MoveIt2", "Dora-MoveIt2"],
    [
        ["安装时间", "30分钟+", "5分钟"],
        ["操作系统", "Ubuntu 22.04 only", "macOS / Linux 均可"],
        ["构建系统", "colcon + CMake", "pip install"],
        ["配置方式", "Setup Assistant GUI", "Python类"],
        ["仿真", "Gazebo + RViz", "MuJoCo（自带可视化）"],
    ])

# 1 - Installation
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "三步完成环境搭建", font_size=28, color=TEXT_WHITE, bold=True)
add_code_block(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(2.0),
"""# Step 1: 安装dora框架
pip install dora-rs

# Step 2: 克隆并安装
git clone https://github.com/dora-rs/dora-moveit2.git
cd dora-moveit2
pip install -e dora_moveit/
pip install -e examples/move_group_demo/
pip install -e examples/dual_gen72/

# Step 3: 安装MuJoCo
pip install mujoco""", font_size=13)

add_text_box(slide, Inches(6.8), Inches(0.4), Inches(5.5), Inches(0.6),
             "环境测试", font_size=28, color=TEXT_WHITE, bold=True)
add_code_block(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(1.5),
"""cd examples/move_group_demo
dora up
dora start dataflows/moveit_example_mujoco.yml
# 看到MuJoCo窗口+机械臂运动 = 成功!
dora stop""", font_size=13)

add_bullet_list(slide, Inches(6.8), Inches(3.0), Inches(5.5), Inches(2.0), [
    "无需colcon、CMake、rosdep",
    "不限操作系统版本（macOS/Linux均可）",
    "Python 3.9+ 即可",
], font_size=14)


# ══════════════════════════════════════════════════
#  CHAPTER 2: Dataflow Basics
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_chapter_header(slide, 2, "Dora 数据流基础", "Dora Dataflow Basics")

# concept mapping
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "ROS2 vs Dora 概念映射", font_size=28, color=TEXT_WHITE, bold=True)
add_table_slide(slide, Inches(0.8), Inches(1.2),
    [Inches(3.0), Inches(3.5), Inches(5.0)],
    ["ROS2 概念", "Dora 等价物", "说明"],
    [
        ["Node（节点）", "Operator（操作器）", "每个Python脚本就是一个操作器"],
        ["Topic（话题）", "数据流连接", "在YAML里声明inputs/outputs"],
        ["Service（服务）", "输出→输入配对", "用配对的输入输出模拟"],
        ["Parameter", "环境变量 (env)", "YAML里的env字段"],
        ["Launch file", "Dataflow YAML", "一个文件定义整个系统"],
        ["TF2 坐标树", "Config中的FK链", "静态配置，无需动态广播"],
        ["URDF/SRDF", "Python Config类", "关节限制、碰撞几何全在Python里"],
    ])

# operator pattern
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "操作器基本模式", font_size=28, color=TEXT_WHITE, bold=True)
add_code_block(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(4.5),
"""from dora import Node
import pyarrow as pa

def main():
    node = Node()
    for event in node:
        if event["type"] == "INPUT":
            if event["id"] == "my_input":
                data = event["value"].to_numpy()
                result = process(data)
                node.send_output(
                    "my_output",
                    pa.array(result)
                )
        elif event["type"] == "STOP":
            break""", font_size=14)

add_text_box(slide, Inches(6.8), Inches(0.4), Inches(5.5), Inches(0.6),
             "数据流YAML（系统接线图）", font_size=24, color=TEXT_WHITE, bold=True)
add_code_block(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(4.5),
"""nodes:
  - id: mujoco_sim
    path: .../main.py
    inputs:
      tick: dora/timer/millis/10
      control_input: executor/joint_commands
    outputs:
      - joint_positions

  - id: planner
    path: .../planner.py
    inputs:
      plan_request: user/plan_request
    outputs:
      - trajectory
    env:
      ROBOT_CONFIG_MODULE: "...config.gen72"
""", font_size=13)

# config injection
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "配置注入机制", font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "库操作器不绑定任何特定机器人 — 配置通过环境变量注入", font_size=18, color=ACCENT_BLUE)
add_code_block(slide, Inches(0.8), Inches(1.8), Inches(5.0), Inches(2.0),
"""# YAML 中设置环境变量
env:
  ROBOT_CONFIG_MODULE:
    "move_group_demo.config.gen72"
""", font_size=14)
add_code_block(slide, Inches(6.2), Inches(1.8), Inches(6.0), Inches(2.0),
"""# 操作器中加载配置
from dora_moveit.config import load_config
config = load_config()
print(config.NUM_JOINTS)    # 7
print(config.HOME_CONFIG)   # [0,-0.5,...]
""", font_size=14)
add_bullet_list(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(2.0), [
    "换机器人只需：写新Config类 + 改YAML env变量",
    "所有库操作器（规划器、IK、执行器）自动适配",
    "支持单臂和双臂配置",
], font_size=16)


# ══════════════════════════════════════════════════
#  CHAPTER 3: GEN72 Model
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_chapter_header(slide, 3, "GEN72 机器人模型与配置", "GEN72 Robot Model & Configuration")

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "GEN72 Python Config 类", font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "代替 URDF/SRDF — 可读、可计算、可继承", font_size=16, color=ACCENT_BLUE)
add_code_block(slide, Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.2),
"""class GEN72Config:
    NUM_JOINTS = 7

    JOINT_LOWER_LIMITS = np.array([
      -3.00, -1.83, -3.00, -2.88,
      -3.00, -1.71, -3.00])
    JOINT_UPPER_LIMITS = np.array([
       3.00,  1.83,  3.00,  0.96,
       3.00,  1.78,  3.00])

    LINK_TRANSFORMS = [
      {"xyz":[0,0,0.218], "rpy":[0,0,0]},
      {"xyz":[0,0,0], "rpy":[-1.57,0,0]},
      # ... 共7个关节变换
    ]

    COLLISION_GEOMETRY = [
      ("sphere", [0.055]),  # base
      ("sphere", [0.045]),  # Link1
      # ... 共8个连杆
    ]

    HOME_CONFIG = np.array(
      [0, -0.5, 0, 0, 0, 0.5, 0])

    NAMED_POSES = {
      "home": HOME_CONFIG,
      "zero": np.zeros(7),
    }""", font_size=12)

add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.5),
             "GEN72 机械臂结构", font_size=22, color=TEXT_WHITE, bold=True)
add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0), [
    "7个旋转关节（joint1 ~ joint7）",
    "8个连杆（base_link + Link1~7）",
    "Link7 = 末端法兰（安装夹爪）",
    "每个关节有角度限制",
    "  例: joint4 ∈ [-2.88, 0.96] rad",
    "  （约 -165° ~ 55°，不对称）",
    "",
    "FK链：base → j1 → L1 → j2 → ...",
    "碰撞几何：球形近似每个连杆",
    "命名姿态：home, safe, zero",
], font_size=15, bullet_color=ACCENT_BLUE)


# ══════════════════════════════════════════════════
#  CHAPTER 4: Configuration
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_chapter_header(slide, 4, "Dora-MoveIt2 配置体系", "Dora-MoveIt2 Configuration")

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "配置三层架构", font_size=28, color=TEXT_WHITE, bold=True)

# Layer 1
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(3.6), Inches(2.5), BG_CARD, 0.03)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(3.2), Inches(0.4),
             "第1层：RobotConfig协议", font_size=16, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(3.2), Inches(1.8),
             "库定义的接口协议\n\nNUM_JOINTS: int\nJOINT_LOWER_LIMITS\nLINK_TRANSFORMS\nCOLLISION_GEOMETRY\nHOME_CONFIG\nNAMED_POSES",
             font_size=12, color=TEXT_LIGHT, font_name="Consolas")

# Layer 2
add_shape_bg(slide, Inches(4.8), Inches(1.3), Inches(3.6), Inches(2.5), BG_CARD, 0.03)
add_text_box(slide, Inches(5.0), Inches(1.4), Inches(3.2), Inches(0.4),
             "第2层：具体Config类", font_size=16, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(5.0), Inches(1.8), Inches(3.2), Inches(1.8),
             "用户编写\n\nclass GEN72Config:\n  NUM_JOINTS = 7\n  JOINT_LOWER_LIMITS\n    = np.array([...])\n  ...",
             font_size=12, color=TEXT_LIGHT, font_name="Consolas")

# Layer 3
add_shape_bg(slide, Inches(8.8), Inches(1.3), Inches(3.6), Inches(2.5), BG_CARD, 0.03)
add_text_box(slide, Inches(9.0), Inches(1.4), Inches(3.2), Inches(0.4),
             "第3层：YAML注入", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(9.0), Inches(1.8), Inches(3.2), Inches(1.8),
             "运行时绑定\n\nenv:\n  ROBOT_CONFIG_MODULE:\n    \"app.config.gen72\"\n\nload_config()自动加载",
             font_size=12, color=TEXT_LIGHT, font_name="Consolas")

# DualArmConfig
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
             "双臂扩展：DualArmConfig", font_size=22, color=ACCENT_ORANGE, bold=True)
add_code_block(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.2),
"""class DualGEN72Config:
    ARM_CHAINS = ["left_arm", "right_arm"]          # 臂链标识
    NUM_JOINTS = 7                                    # 每个臂
    ARM_BASE_TRANSFORMS = {                           # 基座位置
        "left_arm":  {"xyz": [-0.3, 0, 0.8], "rpy": [0, 0, 0]},
        "right_arm": {"xyz": [0.3, 0, 0.8],  "rpy": [0, 0, 3.14]},
    }
    HOME_CONFIG_PER_CHAIN = {"left_arm": HOME, "right_arm": HOME}  # 每臂home""", font_size=12)


# ══════════════════════════════════════════════════
#  CHAPTER 5: MoveGroup API
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_chapter_header(slide, 5, "MoveGroup Python API", "MoveGroup Python API Basics")

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "MoveGroup 核心用法", font_size=28, color=TEXT_WHITE, bold=True)
add_code_block(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.5),
"""from dora_moveit.workflow.move_group \\
    import MoveGroup

group = MoveGroup("gen72")
scene = group.get_planning_scene_interface()

# 1. 命名姿态
group.set_named_target("home")
group.go(wait=True)

# 2. 关节目标
group.go([1.57, -0.785, 0, -1.57,
          0, 0.785, 0], wait=True)

# 3. 末端位姿目标 (IK自动求解)
group.set_pose_target(
    [0.15, 0.1, 0.6, 0, 0, 0])
group.go(wait=True)

# 4. 笛卡尔路径 (直线运动)
traj, frac = group.compute_cartesian_path(
    waypoints, eef_step=0.01)
group.execute(traj, wait=True)

# 5. 避障
scene.add_box("obs", [0,0,0.5],
              [0.1, 0.1, 0.5])
group.go(wait=True)  # 自动避开
scene.remove_world_object("obs")
""", font_size=12)

add_text_box(slide, Inches(7.0), Inches(0.4), Inches(5.5), Inches(0.6),
             "API 对照表", font_size=22, color=TEXT_WHITE, bold=True)
add_table_slide(slide, Inches(7.0), Inches(1.2),
    [Inches(2.8), Inches(2.7)],
    ["功能", "调用方式"],
    [
        ["命名目标", "set_named_target('home')"],
        ["关节目标", "go([j1, j2, ...])"],
        ["位姿目标", "set_pose_target([x,y,z,r,p,y])"],
        ["规划", "plan()"],
        ["执行", "execute(traj)"],
        ["笛卡尔路径", "compute_cartesian_path()"],
        ["停止", "stop()"],
        ["正运动学", "compute_fk()"],
        ["逆运动学", "compute_ik(pose)"],
        ["添加障碍", "scene.add_box(...)"],
    ])


# ══════════════════════════════════════════════════
#  CHAPTER 6: Single Arm Planning
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_chapter_header(slide, 6, "单臂运动规划", "Single Arm Motion Planning")

# joint space
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "关节空间规划 & 笛卡尔路径", font_size=28, color=TEXT_WHITE, bold=True)
add_code_block(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(3.0),
"""# 关节空间规划
group.set_named_target("home")
group.go(wait=True)

joint_goal = group.get_current_joint_values()
joint_goal[0] = 1.57   # 底座旋转90°
joint_goal[1] = -0.785  # 肩关节倾斜
joint_goal[3] = -1.57   # 肘关节弯曲
group.go(joint_goal, wait=True)""", font_size=13)

add_code_block(slide, Inches(0.8), Inches(4.5), Inches(5.8), Inches(2.6),
"""# 笛卡尔路径（直线运动）
pos, _ = group.get_current_pose()
waypoints = [
  [pos[0], pos[1], pos[2]-0.1, 0,0,0],
  [pos[0], pos[1]+0.1, pos[2]-0.1, 0,0,0],
]
traj, frac = group.compute_cartesian_path(
    waypoints, eef_step=0.01)
group.execute(traj, wait=True)""", font_size=13)

add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
             "四大核心操作", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0), [
    "关节空间规划",
    "  直接控制每个关节角度",
    "  最简单，适合精确姿态控制",
    "",
    "笛卡尔路径（直线运动）",
    "  末端沿直线运动",
    "  每1cm插值+IK求解",
    "",
    "简单抓取流程",
    "  接近→下降→抓取→提升→移动→放置",
    "",
    "避障",
    "  添加碰撞物体到PlanningScene",
    "  RRT-Connect自动规划避障路径",
], font_size=14, bullet_color=ACCENT_GREEN)


# ══════════════════════════════════════════════════
#  CHAPTER 7: Dual Arm Coordination
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_chapter_header(slide, 7, "双臂协调控制", "Dual-Arm Coordination")

# architecture
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "14D 统一规划架构", font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "两个7D手臂拼接成14D空间 — 一个规划器同时规划 — 天然保证同步和避碰",
             font_size=16, color=ACCENT_BLUE)

add_code_block(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.5),
"""# 系统数据流
用户 (DualMoveGroup API)
  → plan_request (14D: left_7D + right_7D) → 规划器 (14D RRT-Connect)
  → trajectory (14D) → 执行器 (拆分: left_7D + right_7D)
  → joint_commands (14D) → MuJoCo (14个执行器: 7 left + 7 right)
  → joint_positions (14D) → 规划场景 (per-chain状态跟踪)""", font_size=13)

add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.5),
             "启动双臂环境", font_size=20, color=TEXT_WHITE, bold=True)
add_code_block(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(1.5),
"""cd examples/dual_gen72
dora up
dora start dataflows/dual_gen72_mujoco.yml
# 看到两个GEN72臂 + 桌子 + 球 + 盘子""", font_size=13)

add_text_box(slide, Inches(6.8), Inches(4.3), Inches(5.5), Inches(0.5),
             "MuJoCo 场景", font_size=20, color=TEXT_WHITE, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(4.9), Inches(5.5), Inches(1.5), [
    "左臂（蓝色）x = -0.3",
    "右臂（红色）x = +0.3",
    "红球 (0, 0.1, 0.83) + 绿盘 (0, -0.1, 0.82)",
], font_size=14)

# DualMoveGroup API
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "DualMoveGroup API — 三种运动模式", font_size=28, color=TEXT_WHITE, bold=True)

# mode 1
add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(4.0), Inches(3.2), BG_CARD, 0.03)
add_text_box(slide, Inches(0.7), Inches(1.3), Inches(3.6), Inches(0.4),
             "同步运动", font_size=18, color=ACCENT_GREEN, bold=True)
add_code_block(slide, Inches(0.7), Inches(1.8), Inches(3.6), Inches(2.4),
"""# 双臂同时回home
group.set_named_target(
  left_name="home",
  right_name="home")
group.go(wait=True)

# 指定关节目标
group.go(
  left_joints=[0.5,...],
  right_joints=[-0.5,...],
  wait=True)""", font_size=11)

# mode 2
add_shape_bg(slide, Inches(4.7), Inches(1.2), Inches(4.0), Inches(3.2), BG_CARD, 0.03)
add_text_box(slide, Inches(4.9), Inches(1.3), Inches(3.6), Inches(0.4),
             "独立运动", font_size=18, color=ACCENT_ORANGE, bold=True)
add_code_block(slide, Inches(4.9), Inches(1.8), Inches(3.6), Inches(2.4),
"""# 只动左臂，右臂不动
group.go_left(
  [0, 0.3, 0, -1, 0, 0.5, 0],
  wait=True)

# 只动右臂，左臂不动
group.go_right(
  [0, 0.3, 0, -1, 0, 0.5, 0],
  wait=True)""", font_size=11)

# mode 3
add_shape_bg(slide, Inches(8.9), Inches(1.2), Inches(4.0), Inches(3.2), BG_CARD, 0.03)
add_text_box(slide, Inches(9.1), Inches(1.3), Inches(3.6), Inches(0.4),
             "命名姿态", font_size=18, color=ACCENT_BLUE, bold=True)
add_code_block(slide, Inches(9.1), Inches(1.8), Inches(3.6), Inches(2.4),
"""# 分别设置
group.set_named_target(
  left_name="home")
group.set_named_target(
  right_name="zero")
group.go(wait=True)

# 同时设置
group.set_named_target(
  left_name="home",
  right_name="home")""", font_size=11)

add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(2.0), [
    "14D统一规划：双臂同时搜索路径，天然保证同步到达",
    "自动避碰：check_inter_arm_collision() 检查左臂×右臂所有链接对",
    "向后兼容：单臂MoveGroup API完全不受影响",
], font_size=16, bullet_color=ACCENT_GREEN)


# ══════════════════════════════════════════════════
#  CHAPTER 8: Final Demo
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_chapter_header(slide, 8, "综合实验与调试", "Final Demo & Debug")

# demo flow
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "综合Demo：双臂抓取-传递-放置", font_size=28, color=TEXT_WHITE, bold=True)

steps = [
    ("Step 1", "双臂 → home", "同步", ACCENT_BLUE),
    ("Step 2", "左臂 → 物体上方", "独立", ACCENT_GREEN),
    ("Step 3", "左臂下降抓取", "独立", ACCENT_GREEN),
    ("Step 4", "左臂提升", "独立", ACCENT_GREEN),
    ("Step 5", "右臂 → 接收位置", "独立", ACCENT_ORANGE),
    ("Step 6", "双臂 → 交接位置", "同步", ACCENT_BLUE),
    ("Step 7", "右臂 → 放置位置", "独立", ACCENT_ORANGE),
    ("Step 8", "右臂下降放置", "独立", ACCENT_ORANGE),
    ("Step 9", "右臂撤回", "独立", ACCENT_ORANGE),
    ("Step 10", "双臂 → home", "同步", ACCENT_BLUE),
]
for i, (step, desc, mode, color) in enumerate(steps):
    row = i // 5
    col = i % 5
    x = Inches(0.5 + col * 2.5)
    y = Inches(1.2 + row * 2.8)
    shape = add_shape_bg(slide, x, y, Inches(2.3), Inches(2.2), BG_CARD, 0.05)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.1), Inches(0.35),
                 step, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.5), Inches(2.1), Inches(1.0),
                 desc, font_size=13, color=TEXT_WHITE)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.6), Inches(2.1), Inches(0.35),
                 f"[{mode}]", font_size=11, color=TEXT_DIM)


# debugging
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "常见错误与调试", font_size=28, color=TEXT_WHITE, bold=True)

# error 1
add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(3.9), Inches(3.0), BG_CARD, 0.03)
add_text_box(slide, Inches(0.7), Inches(1.3), Inches(3.5), Inches(0.4),
             "规划失败", font_size=18, color=ACCENT_RED, bold=True)
add_text_box(slide, Inches(0.7), Inches(1.8), Inches(3.5), Inches(0.3),
             "go() 返回 False", font_size=12, color=TEXT_DIM)
add_bullet_list(slide, Inches(0.7), Inches(2.2), Inches(3.5), Inches(1.8), [
    "关节角度超出限制",
    "碰撞/自碰撞",
    "规划超时（路径太复杂）",
    "→ 检查JOINT_LOWER/UPPER_LIMITS",
], font_size=12, bullet_color=ACCENT_RED)

# error 2
add_shape_bg(slide, Inches(4.7), Inches(1.2), Inches(3.9), Inches(3.0), BG_CARD, 0.03)
add_text_box(slide, Inches(4.9), Inches(1.3), Inches(3.5), Inches(0.4),
             "IK求解失败", font_size=18, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(4.9), Inches(1.8), Inches(3.5), Inches(0.3),
             "compute_ik() 返回 None", font_size=12, color=TEXT_DIM)
add_bullet_list(slide, Inches(4.9), Inches(2.2), Inches(3.5), Inches(1.8), [
    "目标位置超出工作空间",
    "姿态不可达",
    "→ 用compute_fk()确认范围",
    "→ 调整目标位置",
], font_size=12, bullet_color=ACCENT_ORANGE)

# error 3
add_shape_bg(slide, Inches(8.9), Inches(1.2), Inches(3.9), Inches(3.0), BG_CARD, 0.03)
add_text_box(slide, Inches(9.1), Inches(1.3), Inches(3.5), Inches(0.4),
             "执行超时", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(9.1), Inches(1.8), Inches(3.5), Inches(0.3),
             "卡在poll_until()", font_size=12, color=TEXT_DIM)
add_bullet_list(slide, Inches(9.1), Inches(2.2), Inches(3.5), Inches(1.8), [
    "轨迹太长/仿真太慢",
    "数据流连接断开",
    "→ dora list 检查节点状态",
    "→ dora stop && dora up 重启",
], font_size=12, bullet_color=ACCENT_BLUE)

# architecture diagram
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
             "系统架构一览", font_size=20, color=TEXT_WHITE, bold=True)
add_code_block(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.0),
"""┌─ 用户代码 (MoveGroup / DualMoveGroup) ──────────────────────────────────┐
│  set_named_target / go / compute_cartesian_path                          │
└────── plan_request ──┬───── ik_request ──┐                               │
       ┌───────────────▼─┐    ┌────────────▼──┐                            │
       │  Planner (RRT)  │    │  IK (TracIK)  │   所有模块通过YAML数据流连接 │
       └───────┬─────────┘    └───────┬───────┘   换机器人只改Config+模型     │
           trajectory            ik_solution                                │
       ┌───────▼──────────────────────┘        ┌─ Planning Scene ─┐        │
       │  Trajectory Executor (插值)           │ (场景管理/碰撞)   │        │
       └───────┬──────────┘                    └──────────────────┘        │
          joint_commands → MuJoCo (仿真+可视化)                              │""", font_size=11)


# ── SUMMARY SLIDE ──
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "课程总结", font_size=36, color=TEXT_WHITE, bold=True)

add_table_slide(slide, Inches(0.5), Inches(1.3),
    [Inches(1.2), Inches(3.5), Inches(4.0), Inches(3.5)],
    ["章节", "主题", "关键代码/命令", "核心收获"],
    [
        ["Ch1", "环境搭建", "pip install -e dora_moveit/", "5分钟装好"],
        ["Ch2", "数据流基础", "YAML + Node事件循环", "理解系统架构"],
        ["Ch3", "GEN72模型", "GEN72Config类", "替代URDF/SRDF"],
        ["Ch4", "配置体系", "RobotConfig协议", "3层配置注入"],
        ["Ch5", "MoveGroup", "go() / set_pose_target()", "API兼容ROS"],
        ["Ch6", "单臂规划", "compute_cartesian_path()", "4大操作"],
        ["Ch7", "双臂协调", "DualMoveGroup / go_left()", "14D统一规划"],
        ["Ch8", "综合实验", "抓取-传递-放置 Demo", "完整流程"],
    ])

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.8),
             "核心优势：纯Python · pip安装 · API兼容ROS MoveIt · 14D双臂统一规划 · MuJoCo仿真",
             font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


# ── THANK YOU ──
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "感谢学习！", font_size=48, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "Dora-MoveIt2 + GEN72 双臂机器人精简课程",
             font_size=20, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "GitHub: github.com/dora-rs/dora-moveit2",
             font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "课后多实操 · 先练单臂再练双臂 · 有问题GitHub提issue",
             font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════
output_path = "/Users/nupylot/Public/github_dora_nav_moveit/dora-moveit2/course-script/DoraMoveIt2_课件.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
